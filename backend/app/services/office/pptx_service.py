import logging

logger = logging.getLogger(__name__)


class PptxService:

    MAX_TEXT_LENGTH = 5000

    async def export(self, file_path: str, json_content: dict) -> None:
        from pptx import Presentation as NewPresentation
        prs = NewPresentation()
        content = json_content.get("content", json_content) if isinstance(json_content, dict) else json_content

        if isinstance(content, list):
            for slide_data in content:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                elements = slide_data.get("elements", [])
                for elem in elements:
                    if elem.get("type") == "textbox":
                        from pptx.util import Inches, Pt
                        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
                        tf = txBox.text_frame
                        tf.text = elem.get("content", "")[:self.MAX_TEXT_LENGTH]

        prs.save(file_path)

    def preview_patch(self, patch: dict, json_content: dict) -> dict:
        if patch.get("operation_type") not in ("replace_text",):
            raise ValueError("PPTX 补丁仅支持 replace_text 操作类型")
        return {"preview_passed": True, "risk_level": "medium"}
