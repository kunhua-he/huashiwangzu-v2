from app.middleware.auth import require_permission
from app.models.user import User
from app.schemas.common import ApiResponse
from app.services.module_registry import register_capability
from app.services.uploaded_file_runner import run_uploaded_file_capability
from fastapi import APIRouter, Depends
from pydantic import BaseModel

router = APIRouter(prefix="/api/pptx-parser", tags=["pptx-parser"])


class ParseRequest(BaseModel):
    file_id: int


async def _parse(params: dict, caller: str) -> dict:
    from pptx import Presentation
    import base64, io, hashlib
    from app.services.module_registry import call_capability

    allowed = {"pptx"}

    def parse_file(file_id, _file, full_path, _ext):
        prs = Presentation(str(full_path))
        blocks = []
        resources = []
        resource_counter = 0
        resource_map: dict[int, dict] = {}

        for slide_idx, slide in enumerate(prs.slides):
            pno = slide_idx + 1
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        block_type = "heading" if ("title" in str(shape.name).lower() or "标题" in str(shape.name)) else "paragraph"
                        blocks.append({"type": block_type, "text": text, "page": pno, "resource_ref": None})
                if shape.shape_type and "picture" in str(shape.shape_type).lower():
                    resource_counter += 1
                    try:
                        img = shape.image
                        img_bytes = img.blob
                        resource_map[resource_counter] = img_bytes
                    except Exception:
                        img_bytes = b""

                    blocks.append({"type": "image", "text": "", "page": pno, "resource_ref": resource_counter})
                    resources.append({
                        "id": resource_counter,
                        "type": "image",
                        "mime_type": shape.image.content_type if hasattr(shape, "image") and hasattr(shape.image, "content_type") else "image/png",
                        "filename": f"slide{pno}_{hashlib.md5(str(shape.name).encode()).hexdigest()[:8]}.png",
                        "description": f"Slide {pno} image ({shape.name})",
                        "_bytes_b64": base64.b64encode(img_bytes).decode("ascii") if img_bytes else "",
                    })

        return {
            "file_id": file_id,
            "format": "pptx",
            "blocks": blocks,
            "resources": resources,
        }

    result = await run_uploaded_file_capability(params, caller, allowed, parse_file)

    # Store extracted resources via content:store_resource capability
    for res in result.get("resources", []):
        data_b64 = res.pop("_bytes_b64", "")
        if data_b64:
            try:
                await call_capability(
                    "content", "store_resource",
                    {
                        "data_b64": data_b64,
                        "resource_type": "image",
                        "mime_type": res.get("mime_type", "image/png"),
                        "filename": res.get("filename", "resource.png"),
                        "description": res.get("description", ""),
                    },
                    caller,
                )
            except Exception:
                pass

    return result


@router.get("/health")
async def health():
    return ApiResponse(data={"module": "pptx-parser", "status": "ok"})


@router.post("/parse")
async def call_parse(payload: ParseRequest, user: User = Depends(require_permission("viewer"))):
    result = await _parse({"file_id": payload.file_id}, f"user:{user.id}")
    return ApiResponse(data=result)


register_capability(
    "pptx-parser", "parse", _parse,
    description="Parse PPTX files into unified content blocks",
    brief="解析 PPTX 文档",
    parameters={"file_id": {"type": "int"}},
    min_role="viewer",
)
