
"""Sandbox test for pptx-parser module."""
import sys
from pathlib import Path
SAMPLE = Path(__file__).resolve().parent / "samples" / "sample.pptx"
if not SAMPLE.exists():
    print("ERROR: sample.pptx not found"); sys.exit(1)
from pptx import Presentation


def parse_pptx(path):
    prs = Presentation(str(path))
    blocks, resources, rc = [], [], 0
    for si, slide in enumerate(prs.slides):
        pno = si + 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if not t: continue
                    bt = "标题" if "title" in str(shape.name).lower() else "段落"
                    blocks.append({"type": bt, "text": t, "page": pno, "resource_ref": None})
            if shape.shape_type and "picture" in str(shape.shape_type).lower():
                rc += 1
                blocks.append({"type": "图片", "text": "", "page": pno, "resource_ref": rc})
                resources.append({"id": rc, "type": "图片", "file_storage_id": None, "text_desc": f"Slide {pno} image ({shape.name})"})
    return {"file_id": 0, "format": "pptx", "blocks": blocks, "resources": resources}


def validate(result):
    assert result["format"] == "pptx"
    for b in result["blocks"]: assert all(k in b for k in ("type","text","page","resource_ref"))
    print("  Validation PASS (%d blocks, %d resources)" % (len(result["blocks"]), len(result["resources"])))


def main():
    print("="*60); print("pptx-parser sandbox test"); print("="*60)
    result = parse_pptx(SAMPLE)
    for b in result["blocks"]:
        print("    [%s] p=%s %s" % (b["type"], b["page"], b["text"][:70]))
    validate(result)
    print("PASS: pptx-parser sandbox test")

if __name__ == "__main__": main()
