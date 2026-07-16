from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SOURCE = Path("/Users/hekunhua/Documents/Hermes工作区/竞品分析/直播文案_01版.md")
OUTPUT = Path("/Users/hekunhua/Documents/Hermes工作区/竞品分析/全肌品牌招商直播_精校原稿PPT.pptx")

W = Inches(13.333)
H = Inches(7.5)

BG = "F7F4EF"
INK = "1D1D1F"
MUTED = "6E6B66"
LINE = "D8D2C8"
WHITE = "FFFFFF"
ORANGE = "F15A3A"
YELLOW = "F4C552"
GREEN = "28776D"
BLUE = "3778A8"
PALE_ORANGE = "F8DED5"
PALE_GREEN = "DCEAE6"
PALE_BLUE = "DFEAF2"
PALE_YELLOW = "F7EBC8"

FONT_CN = "PingFang SC"
FONT_NUM = "Avenir Next"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def extract(source: str, start: str, end: str | None = None) -> str:
    start_pos = source.find(start)
    if start_pos < 0:
        return ""
    end_pos = source.find(end, start_pos + len(start)) if end else -1
    if end_pos < 0:
        end_pos = len(source)
    return source[start_pos:end_pos].strip()


def set_run(run, size: float, color: str, bold: bool = False, font: str = FONT_CN):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
    font: str = FONT_CN,
):
    # Allow concise positional calls where alignment follows color directly.
    if not isinstance(bold, (bool, type(None))):
        align = bold
        bold = False
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = 1.08
    r = p.add_run()
    r.text = text
    set_run(r, size, color, bold, font)
    return box


def add_rich_text(
    slide,
    parts: list[tuple[str, float, str, bool]],
    x: float,
    y: float,
    w: float,
    h: float,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    for text, size, color, bold in parts:
        r = p.add_run()
        r.text = text
        set_run(r, size, color, bold, FONT_NUM if any(ch.isdigit() for ch in text) else FONT_CN)
    return box


def add_rect(slide, x: float, y: float, w: float, h: float, fill: str, line: str | None = None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def add_circle(slide, x: float, y: float, d: float, fill: str, line: str | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: str = LINE, width: float = 1.5):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_arrow(slide, x: float, y: float, w: float, h: float, fill: str):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(fill)
    return shape


def add_bullets(
    slide,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 17,
    color: str = INK,
    accent: str = ORANGE,
    gap: float = 7,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        p.level = 0
        r = p.add_run()
        r.text = "●  "
        set_run(r, size - 2, accent, True)
        r = p.add_run()
        r.text = item
        set_run(r, size, color, False)
    return box


def add_background(slide, color: str = BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(color)


def add_header(slide, section: str, title: str, subtitle: str | None = None, page: int | None = None):
    add_rect(slide, 0.62, 0.56, 0.08, 0.48, ORANGE)
    add_text(slide, section.upper(), 0.83, 0.52, 2.5, 0.35, 10, ORANGE, True, font=FONT_NUM)
    add_text(slide, title, 0.81, 0.88, 11.4, 0.58, 27, INK, True)
    if subtitle:
        add_text(slide, subtitle, 0.84, 1.48, 11.2, 0.4, 12.5, MUTED)
    add_line(slide, 0.82, 1.93, 12.55, 1.93, LINE, 1)
    if page is not None:
        add_text(slide, f"{page:02d}", 12.26, 0.54, 0.45, 0.28, 9.5, MUTED, True, PP_ALIGN.RIGHT, font=FONT_NUM)


def add_footer(slide, page: int):
    add_line(slide, 0.82, 7.07, 12.55, 7.07, LINE, 0.7)
    add_text(slide, "全肌品牌招商直播｜精校原稿PPT", 0.82, 7.12, 4.2, 0.2, 8.5, MUTED)
    add_text(slide, f"{page:02d}", 12.18, 7.12, 0.36, 0.2, 8.5, MUTED, True, PP_ALIGN.RIGHT, font=FONT_NUM)


def add_note(slide, note: str):
    tf = slide.notes_slide.notes_text_frame
    tf.text = note


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    fill: str = WHITE,
    accent: str = ORANGE,
    title_size: float = 17,
    body_size: float = 13.5,
):
    add_rect(slide, x, y, w, h, fill, LINE)
    add_rect(slide, x, y, 0.08, h, accent)
    add_text(slide, title, x + 0.25, y + 0.18, w - 0.42, 0.42, title_size, INK, True)
    add_text(slide, body, x + 0.25, y + 0.68, w - 0.42, h - 0.82, body_size, MUTED)


def new_slide(prs: Presentation, page: int, section: str, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, section, title, subtitle, page)
    add_footer(slide, page)
    return slide


def build_deck():
    source = SOURCE.read_text(encoding="utf-8")
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    prs.core_properties.title = "全肌品牌招商直播｜精校原稿PPT"
    prs.core_properties.subject = "竞品招商直播内容还原与培训演示"
    prs.core_properties.author = "基于用户提供的ASR精校稿生成"

    # 01 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, INK)
    add_rect(slide, 0, 0, 0.22, 7.5, ORANGE)
    add_rect(slide, 9.55, 0, 3.78, 7.5, ORANGE)
    for i, label in enumerate(["经营底盘", "获客裂变", "三张卡", "陪跑对赌"]):
        y = 1.05 + i * 1.18
        add_text(slide, f"0{i+1}", 9.9, y, 0.55, 0.3, 10, WHITE, True, font=FONT_NUM)
        add_text(slide, label, 10.5, y - 0.08, 2.1, 0.42, 18, WHITE, True)
        add_line(slide, 9.9, y + 0.46, 12.65, y + 0.46, "F58B72", 0.8)
    add_text(slide, "竞品原始干净稿 · 演示版", 0.85, 0.82, 5.2, 0.36, 12, YELLOW, True)
    add_text(slide, "全肌品牌\n招商直播", 0.82, 1.55, 7.8, 1.9, 38, WHITE, True)
    add_text(slide, "从门店焦虑到3980元对赌成交的完整逻辑", 0.87, 3.72, 7.8, 0.58, 19, "D9D6D0")
    add_text(slide, "精校原稿PPT｜基于三段直播录屏整理", 0.88, 6.42, 6.4, 0.3, 11, "A7A39D")
    add_note(slide, extract(source, "# 全肌品牌招商直播", "## 第一部分"))

    # 02 Reading guide
    page = 2
    slide = new_slide(prs, page, "使用说明", "这是一套“还原稿”，不是我方方案")
    add_card(slide, 0.85, 2.25, 3.72, 3.85, "保留", "讲解顺序\n核心原话\n案例与互动\n成交推进\n异议处理", PALE_GREEN, GREEN, 19, 17)
    add_card(slide, 4.8, 2.25, 3.72, 3.85, "清理", "口头重复\n设备杂音\n视频乱码\n无法辨认插话\nASR明显错字", PALE_YELLOW, YELLOW, 19, 17)
    add_card(slide, 8.75, 2.25, 3.72, 3.85, "不处理", "不改成俏小喵\n不替竞品优化政策\n不验证效果宣称\n不美化收益数据\n不混入二次分析", PALE_ORANGE, ORANGE, 19, 17)
    add_note(slide, extract(source, "> 本稿依据", "---"))

    # 03 Full logic map
    page += 1
    slide = new_slide(prs, page, "全场地图", "一场直播，九步完成认知教育与合作成交", "每一步都在回答门店老板的一个核心问题")
    labels = [
        ("01", "经营问题", "为什么越来越难做"),
        ("02", "价值承诺", "六个月得到什么"),
        ("03", "两条路径", "新客从哪里来"),
        ("04", "980裂变", "为什么客户会成交"),
        ("05", "三张卡", "如何持续复购"),
        ("06", "员工激励", "谁来执行"),
        ("07", "合作方式", "门店怎么加入"),
        ("08", "Q&A", "风险和异议怎么解"),
        ("09", "成交收口", "下一步做什么"),
    ]
    for i, (num, title, body) in enumerate(labels):
        row, col = divmod(i, 3)
        x, y = 0.86 + col * 4.05, 2.22 + row * 1.43
        fill = [PALE_ORANGE, PALE_GREEN, PALE_BLUE][row]
        accent = [ORANGE, GREEN, BLUE][row]
        add_rect(slide, x, y, 3.72, 1.12, fill, fill)
        add_text(slide, num, x + 0.18, y + 0.18, 0.55, 0.28, 10.5, accent, True, font=FONT_NUM)
        add_text(slide, title, x + 0.82, y + 0.13, 2.6, 0.33, 17, INK, True)
        add_text(slide, body, x + 0.82, y + 0.55, 2.6, 0.3, 12.5, MUTED)
    add_note(slide, "PPT结构索引。完整原话见各页演讲者备注。")

    # 04 Opening problem
    page += 1
    slide = new_slide(prs, page, "第一部分", "先把所有焦虑统一成一个问题", "缺客、流失、低到店、低复购、团队无自销能力")
    pains = ["新客进不来", "老客在流失", "到店率太差", "活动越做越疲", "员工不会自销"]
    for i, pain in enumerate(pains):
        x = 0.9 + i * 2.44
        add_circle(slide, x + 0.49, 2.45, 1.26, [ORANGE, YELLOW, GREEN, BLUE, INK][i])
        add_text(slide, f"0{i+1}", x + 0.78, 2.82, 0.68, 0.3, 12, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
        add_text(slide, pain, x, 3.95, 2.25, 0.42, 16, INK, True, PP_ALIGN.CENTER)
    add_rect(slide, 2.0, 5.15, 9.3, 0.88, INK)
    add_text(slide, "它们不是五个零散问题，而是同一个问题：门店没有经营底盘", 2.25, 5.39, 8.8, 0.38, 19, WHITE, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 开场：", "### 五项核心数据指标"))

    # 05 Bottom line
    page += 1
    slide = new_slide(prs, page, "经营底盘", "底盘不是一场业绩，而是一种业绩结构")
    add_text(slide, "80%", 0.9, 2.25, 2.4, 0.85, 44, ORANGE, True, font=FONT_NUM)
    add_text(slide, "业绩", 3.02, 2.55, 1.0, 0.4, 17, MUTED, True)
    add_text(slide, "来自", 4.05, 2.55, 0.8, 0.4, 15, MUTED)
    add_text(slide, "20%", 4.95, 2.25, 2.4, 0.85, 44, GREEN, True, font=FONT_NUM)
    add_text(slide, "客户", 7.08, 2.55, 1.0, 0.4, 17, MUTED, True)
    add_line(slide, 0.95, 3.38, 8.05, 3.38, LINE, 1.3)
    add_bullets(slide, ["没业绩就上新项目、搞活动、请卖手", "反复消耗少数老客，新客仍然进不来", "外部老师一走，团队能力没有留下"], 0.95, 3.75, 7.35, 2.2, 16.5)
    add_rect(slide, 8.8, 2.3, 3.55, 3.58, PALE_YELLOW, PALE_YELLOW)
    add_text(slide, "底盘", 9.15, 2.68, 2.85, 0.56, 27, INK, True, PP_ALIGN.CENTER)
    add_text(slide, "稳定、可持续的\n业绩结构", 9.15, 3.55, 2.85, 1.2, 22, ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, "= 老板经营的底气", 9.15, 5.08, 2.85, 0.4, 14, MUTED, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "经营底盘，就等于", "门店经营管理没有那么难"))

    # 06 Five metrics
    page += 1
    slide = new_slide(prs, page, "五项数据", "门店业绩被拆成五个可管理变量")
    metrics = [
        ("01", "有效客户", "年复购≥3次\n年消费≥1万", ORANGE),
        ("02", "到店率", "目标70%\n约10天1次", YELLOW),
        ("03", "客单次", "一次到店\n≥2个项目", GREEN),
        ("04", "客单价", "按投资回报\n倒推定价", BLUE),
        ("05", "复购率", "约4个月\n刷1张卡", INK),
    ]
    for i, (num, title, body, color) in enumerate(metrics):
        x = 0.86 + i * 2.45
        add_rect(slide, x, 2.35, 2.14, 3.55, WHITE, LINE)
        add_rect(slide, x, 2.35, 2.14, 0.16, color)
        add_text(slide, num, x + 0.2, 2.76, 0.55, 0.3, 11, color, True, font=FONT_NUM)
        add_text(slide, title, x + 0.2, 3.18, 1.74, 0.45, 18, INK, True)
        add_line(slide, x + 0.2, 3.78, x + 1.9, 3.78, LINE, 0.8)
        add_text(slide, body, x + 0.2, 4.14, 1.72, 1.1, 14.5, MUTED, False, PP_ALIGN.LEFT)
    add_text(slide, "业绩 = 有效客户 × 到店率 × 客单次 × 客单价 × 复购率", 1.05, 6.3, 11.2, 0.42, 19, INK, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 五项核心数据指标", "### 门店像工厂"))

    # 07 Effective customer and price
    page += 1
    slide = new_slide(prs, page, "关键诊断", "客单价定高，反而把有效客户挡在门外")
    add_rect(slide, 0.9, 2.35, 4.0, 3.55, PALE_ORANGE, PALE_ORANGE)
    add_text(slide, "案例", 1.2, 2.7, 0.8, 0.34, 11, ORANGE, True)
    add_text(slide, "600㎡", 1.18, 3.12, 2.25, 0.7, 34, INK, True, font=FONT_NUM)
    add_text(slide, "开店10年", 1.2, 3.91, 2.6, 0.4, 18, INK, True)
    add_text(slide, "有效客户只有30–40个", 1.2, 4.57, 3.15, 0.5, 19, ORANGE, True)
    add_text(slide, "客单价却定到500–600元", 1.2, 5.24, 3.2, 0.35, 14, MUTED)
    add_arrow(slide, 5.2, 3.55, 1.35, 0.72, ORANGE)
    add_rect(slide, 6.85, 2.35, 5.35, 3.55, PALE_GREEN, PALE_GREEN)
    add_text(slide, "连锁反应", 7.18, 2.7, 1.2, 0.34, 11, GREEN, True)
    add_bullets(slide, ["获取有效客户周期被无限拉长", "客户到店率下降", "客户复购频次降低"], 7.18, 3.24, 4.45, 1.95, 18, INK, GREEN, 10)
    add_text(slide, "定价公式：门店投资 + 月利润目标 → 倒推合理客单价", 7.18, 5.35, 4.45, 0.42, 14.5, GREEN, True)
    add_note(slide, extract(source, "#### 第一项：有效客户数量", "#### 第五项：复购率"))

    # 08 Bed production line
    page += 1
    slide = new_slide(prs, page, "床位模型", "门店像工厂，床位就是生产线")
    add_rect(slide, 0.95, 2.45, 2.1, 2.7, INK)
    add_text(slide, "30万", 1.2, 2.93, 1.6, 0.68, 34, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "开店投资", 1.2, 3.78, 1.6, 0.4, 15, "D8D4CD", True, PP_ALIGN.CENTER)
    add_arrow(slide, 3.35, 3.48, 1.15, 0.7, ORANGE)
    for i in range(5):
        x = 4.82 + i * 1.35
        add_rect(slide, x, 2.75, 1.0, 2.15, [PALE_ORANGE, PALE_YELLOW, PALE_GREEN, PALE_BLUE, "E6E3DE"][i], LINE)
        add_text(slide, f"床位 {i+1}", x + 0.12, 3.06, 0.76, 0.3, 12, INK, True, PP_ALIGN.CENTER)
        add_text(slide, "3", x + 0.22, 3.58, 0.56, 0.5, 27, [ORANGE, "B58C20", GREEN, BLUE, INK][i], True, PP_ALIGN.CENTER, font=FONT_NUM)
        add_text(slide, "人/天", x + 0.18, 4.2, 0.64, 0.27, 10.5, MUTED, PP_ALIGN.CENTER)
    add_rect(slide, 4.75, 5.55, 6.55, 0.62, PALE_YELLOW, PALE_YELLOW)
    add_text(slide, "任何一张床长期闲置 = 一条生产线没有产出", 5.0, 5.74, 6.0, 0.3, 16.5, INK, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 门店像工厂", "## 第二部分"))

    # 09 Health gap
    page += 1
    slide = new_slide(prs, page, "数据体检", "左边是健康值，右边是门店现状")
    add_rect(slide, 1.0, 2.35, 4.65, 3.5, PALE_GREEN, PALE_GREEN)
    add_text(slide, "健康指标", 1.35, 2.72, 3.95, 0.5, 25, GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, "根据床位数倒推\n客户数 · 到店率 · 客单次\n客单价 · 复购率", 1.4, 3.55, 3.85, 1.65, 18, INK, False, PP_ALIGN.CENTER)
    add_rect(slide, 7.68, 2.35, 4.65, 3.5, PALE_ORANGE, PALE_ORANGE)
    add_text(slide, "现状数据", 8.03, 2.72, 3.95, 0.5, 25, ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, "真实盘点结果\n不是感觉、不是估算\n而是门店经营体检", 8.08, 3.55, 3.85, 1.65, 18, INK, False, PP_ALIGN.CENTER)
    add_arrow(slide, 5.92, 3.55, 1.15, 0.72, INK)
    add_text(slide, "差距", 5.98, 4.42, 1.0, 0.38, 17, ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, "六个月陪跑 = 把差距一项项填回来", 2.1, 6.25, 9.0, 0.42, 20, INK, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "当健康指标算出来", "## 第二部分"))

    # 10 Three hundreds
    page += 1
    slide = new_slide(prs, page, "价值承诺", "六个月“三个一百”")
    cards = [
        ("100", "新增新客", "成交980元\n必须是真新客", ORANGE, PALE_ORANGE),
        ("100", "有效客户", "年复购3次\n年消费1万元以上", GREEN, PALE_GREEN),
        ("100", "家居客户", "能够持续\n月月复购产品", BLUE, PALE_BLUE),
    ]
    for i, (num, title, body, accent, fill) in enumerate(cards):
        x = 0.98 + i * 4.12
        add_rect(slide, x, 2.42, 3.62, 3.42, fill, fill)
        add_text(slide, num, x + 0.35, 2.75, 2.92, 0.9, 46, accent, True, PP_ALIGN.CENTER, font=FONT_NUM)
        add_text(slide, title, x + 0.35, 3.75, 2.92, 0.42, 20, INK, True, PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.35, 4.48, 2.92, 0.92, 16, MUTED, False, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 六个月“三个一百”", "### 过程价值"))

    # 11 Process value
    page += 1
    slide = new_slide(prs, page, "过程价值", "真正要留下来的，是团队的自销能力")
    skills = ["专业知识", "技术手法", "服务能力", "面诊能力", "成交能力", "客户管理"]
    for i, skill in enumerate(skills):
        row, col = divmod(i, 3)
        x, y = 1.1 + col * 3.9, 2.45 + row * 1.46
        add_rect(slide, x, y, 3.35, 1.08, WHITE, LINE)
        add_text(slide, f"0{i+1}", x + 0.22, y + 0.31, 0.48, 0.3, 10.5, ORANGE, True, font=FONT_NUM)
        add_text(slide, skill, x + 0.82, y + 0.25, 2.12, 0.42, 18, INK, True)
    add_rect(slide, 2.0, 5.55, 9.3, 0.72, INK)
    add_text(slide, "老师不是替门店做一场业绩，而是把自销系统留在门店", 2.3, 5.76, 8.7, 0.35, 18.5, WHITE, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 过程价值", "### 公司凭什么"))

    # 12 Two paths
    page += 1
    slide = new_slide(prs, page, "两大路径", "公域引流 + 私域裂变，形成门店流量闭环")
    add_rect(slide, 0.92, 2.45, 5.25, 3.3, PALE_BLUE, PALE_BLUE)
    add_text(slide, "PATH 01", 1.28, 2.8, 1.25, 0.3, 10.5, BLUE, True, font=FONT_NUM)
    add_text(slide, "线上公域", 1.25, 3.28, 2.7, 0.52, 26, INK, True)
    add_bullets(slide, ["抖音 / 美团运营", "80元左右刚需爆品", "月均30+到店新客"], 1.25, 4.05, 4.1, 1.28, 16, INK, BLUE, 6)
    add_rect(slide, 7.16, 2.45, 5.25, 3.3, PALE_GREEN, PALE_GREEN)
    add_text(slide, "PATH 02", 7.52, 2.8, 1.25, 0.3, 10.5, GREEN, True, font=FONT_NUM)
    add_text(slide, "线下私域", 7.49, 3.28, 2.7, 0.52, 26, INK, True)
    add_bullets(slide, ["980元信任卡", "高定家居精华", "A→B→C持续裂变"], 7.49, 4.05, 4.1, 1.28, 16, INK, GREEN, 6)
    add_text(slide, "只做一条腿，走不稳", 4.58, 6.25, 4.2, 0.4, 19, ORANGE, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 公司凭什么", "### 路径一"))

    # 13 Online acquisition
    page += 1
    slide = new_slide(prs, page, "公域引流", "不做9.9元低价团单，只做80元左右刚需爆品")
    add_text(slide, "9.9 / 19.9 / 29.9", 0.95, 2.45, 4.75, 0.62, 28, MUTED, True, font=FONT_NUM)
    add_text(slide, "低价带来流量，也带来低质量客户", 0.98, 3.26, 4.7, 0.42, 16, MUTED)
    add_line(slide, 0.98, 3.9, 5.55, 3.9, LINE, 1)
    add_text(slide, "83元", 0.95, 4.22, 2.7, 0.9, 45, ORANGE, True, font=FONT_NUM)
    add_text(slide, "平均新客成交价", 3.1, 4.55, 2.25, 0.4, 15, INK, True)
    add_rect(slide, 6.35, 2.35, 5.85, 3.7, PALE_ORANGE, PALE_ORANGE)
    add_text(slide, "主推团单", 6.72, 2.74, 1.2, 0.32, 11, ORANGE, True)
    add_text(slide, "毛孔毛囊清洁\n祛痘护理", 6.7, 3.33, 4.75, 1.22, 29, INK, True)
    add_text(slide, "痛点型 · 刚需型 · 高普及 · 高复购", 6.72, 5.12, 4.75, 0.38, 15.5, ORANGE, True)
    add_note(slide, extract(source, "### 路径一：", "#### 为什么团单"))

    # 14 Bun vs seafood
    page += 1
    slide = new_slide(prs, page, "品类选择", "流量大小不等于客户质量")
    add_rect(slide, 0.95, 2.35, 5.25, 3.85, WHITE, LINE)
    add_text(slide, "包子铺", 1.35, 2.78, 2.0, 0.5, 25, INK, True)
    add_text(slide, "大流量", 1.35, 3.58, 1.7, 0.42, 20, ORANGE, True)
    add_text(slide, "低客单", 3.58, 3.58, 1.7, 0.42, 20, MUTED, True)
    add_bullets(slide, ["买包子的人多", "单次消费十几元", "价格导向明显"], 1.35, 4.35, 4.15, 1.32, 16)
    add_rect(slide, 7.12, 2.35, 5.25, 3.85, PALE_GREEN, PALE_GREEN)
    add_text(slide, "海鲜店", 7.52, 2.78, 2.0, 0.5, 25, INK, True)
    add_text(slide, "小流量", 7.52, 3.58, 1.7, 0.42, 20, MUTED, True)
    add_text(slide, "高质量", 9.75, 3.58, 1.7, 0.42, 20, GREEN, True)
    add_bullets(slide, ["人数可能更少", "客单可达几百甚至几千", "问题与价值导向"], 7.52, 4.35, 4.15, 1.32, 16, INK, GREEN)
    add_note(slide, extract(source, "#### 为什么团单", "#### 高价团单"))

    # 15 Video formula
    page += 1
    slide = new_slide(prs, page, "内容支撑", "一分钟短视频，用五个元素证明“凭什么更贵”")
    add_circle(slide, 5.45, 2.72, 2.42, INK)
    add_text(slide, "高价\n团单价值", 5.75, 3.25, 1.82, 0.92, 22, WHITE, True, PP_ALIGN.CENTER)
    elements = [
        ("服务流程", 1.05, 2.35, ORANGE),
        ("使用产品", 1.25, 4.72, YELLOW),
        ("工具呈现", 5.12, 5.48, GREEN),
        ("技术手法", 9.2, 4.72, BLUE),
        ("可视效果", 9.35, 2.35, ORANGE),
    ]
    for idx, (label, x, y, color) in enumerate(elements, 1):
        add_rect(slide, x, y, 2.6, 0.92, WHITE, LINE)
        add_text(slide, f"0{idx}", x + 0.18, y + 0.26, 0.5, 0.28, 10, color, True, font=FONT_NUM)
        add_text(slide, label, x + 0.82, y + 0.2, 1.55, 0.38, 16.5, INK, True)
        add_line(slide, x + (2.6 if x < 5 else 0), y + 0.46, 6.66, 3.93, color, 1.2)
    add_note(slide, extract(source, "#### 高价团单", "#### 不是把人"))

    # 16 Five steps
    page += 1
    slide = new_slide(prs, page, "新客转化", "从进店前到现场裂变：成交五部曲")
    steps = [
        ("01", "进店前", "差异认知", ORANGE),
        ("02", "面诊间", "专业期待", YELLOW),
        ("03", "护理房", "效果验证", GREEN),
        ("04", "抛单", "成交980", BLUE),
        ("05", "回面诊间", "信任+裂变", INK),
    ]
    for i, (num, location, action, color) in enumerate(steps):
        x = 0.84 + i * 2.48
        add_circle(slide, x + 0.55, 2.55, 1.18, color)
        add_text(slide, num, x + 0.79, 2.92, 0.7, 0.28, 11, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
        if i < 4:
            add_arrow(slide, x + 1.83, 2.93, 0.58, 0.38, LINE)
        add_text(slide, location, x, 4.05, 2.25, 0.35, 14, MUTED, True, PP_ALIGN.CENTER)
        add_text(slide, action, x, 4.53, 2.25, 0.45, 18, INK, True, PP_ALIGN.CENTER)
    add_rect(slide, 1.7, 5.63, 9.95, 0.65, PALE_YELLOW, PALE_YELLOW)
    add_text(slide, "每一步都要训练：说什么、做什么、如何进入下一步", 1.95, 5.83, 9.45, 0.3, 16.5, INK, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "#### 新客成交五部曲", "### 路径二"))

    # 17 Private referral
    page += 1
    slide = new_slide(prs, page, "私域裂变", "980元成交后，现场就启动转介绍")
    add_rect(slide, 0.95, 2.38, 3.2, 3.62, PALE_ORANGE, PALE_ORANGE)
    add_text(slide, "980", 1.3, 2.85, 2.5, 0.82, 43, ORANGE, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "新客信任卡", 1.3, 3.77, 2.5, 0.42, 18, INK, True, PP_ALIGN.CENTER)
    add_arrow(slide, 4.48, 3.55, 1.0, 0.68, ORANGE)
    add_rect(slide, 5.82, 2.38, 2.55, 3.62, PALE_GREEN, PALE_GREEN)
    add_text(slide, "2套", 6.15, 2.85, 1.9, 0.72, 35, GREEN, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "高定精华", 6.15, 3.71, 1.9, 0.42, 18, INK, True, PP_ALIGN.CENTER)
    add_text(slide, "1套自用\n1套送朋友", 6.15, 4.45, 1.9, 0.86, 15, MUTED, False, PP_ALIGN.CENTER)
    add_arrow(slide, 8.7, 3.55, 1.0, 0.68, GREEN)
    add_rect(slide, 10.03, 2.38, 2.35, 3.62, INK, INK)
    add_text(slide, "12套", 10.34, 2.85, 1.73, 0.72, 35, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "朋友成交后", 10.34, 3.78, 1.73, 0.42, 16, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "每月到店\n配置1套", 10.34, 4.48, 1.73, 0.84, 14.5, "D8D4CD", False, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 路径二", "#### A、B、C"))

    # 18 ABC chain
    page += 1
    slide = new_slide(prs, page, "裂变链路", "A送B，B成交；B再送C，链路继续")
    people = [
        ("A", "购买980", "送B一套", ORANGE),
        ("B", "到店体验", "成交980", GREEN),
        ("C", "到店体验", "继续裂变", BLUE),
    ]
    for i, (name, line1, line2, color) in enumerate(people):
        x = 1.15 + i * 4.15
        add_circle(slide, x + 0.85, 2.35, 1.55, color)
        add_text(slide, name, x + 1.17, 2.78, 0.9, 0.46, 25, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
        add_text(slide, line1, x + 0.35, 4.18, 2.55, 0.38, 16, INK, True, PP_ALIGN.CENTER)
        add_text(slide, line2, x + 0.35, 4.67, 2.55, 0.35, 14, MUTED, PP_ALIGN.CENTER)
        if i < 2:
            add_arrow(slide, x + 2.85, 2.88, 1.1, 0.58, color)
    add_rect(slide, 3.42, 5.55, 6.5, 0.72, INK)
    add_text(slide, "客户请客，我买单", 3.75, 5.75, 5.85, 0.36, 21, WHITE, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "#### A、B、C裂变链路", "## 第四部分"))

    # 19 Hero project
    page += 1
    slide = new_slide(prs, page, "抓手项目", "提高留客率与到店率，项目必须同时具备三点")
    traits = [
        ("01", "有特色", "即时效果强\n周边形成差异", ORANGE, PALE_ORANGE),
        ("02", "高普及", "绝大多数客户\n都能体验", GREEN, PALE_GREEN),
        ("03", "成瘾性", "10天不做\n就觉得难受", BLUE, PALE_BLUE),
    ]
    for i, (num, title, body, accent, fill) in enumerate(traits):
        x = 1.0 + i * 4.1
        add_rect(slide, x, 2.45, 3.6, 3.35, fill, fill)
        add_text(slide, num, x + 0.32, 2.78, 0.55, 0.3, 11, accent, True, font=FONT_NUM)
        add_text(slide, title, x + 0.32, 3.32, 2.95, 0.5, 24, INK, True)
        add_text(slide, body, x + 0.32, 4.25, 2.95, 0.94, 17, MUTED)
    add_note(slide, extract(source, "### 抓手项目", "### 980元卡的五大"))

    # 20 Five reasons
    page += 1
    slide = new_slide(prs, page, "980成交", "五个理由，把“划算、有效、能送人”叠在一起")
    reasons = [
        ("01", "古法毛囊清洁", "刚需 + 可视化", ORANGE),
        ("02", "水置换", "叠加项目价值", YELLOW),
        ("03", "19个项目", "均价50多元", GREEN),
        ("04", "两套高定精华", "自用 + 送朋友", BLUE),
        ("05", "再得12套", "一年到店抓手", INK),
    ]
    for i, (num, title, body, accent) in enumerate(reasons):
        y = 2.32 + i * 0.78
        add_rect(slide, 1.0, y, 11.25, 0.62, WHITE, LINE)
        add_text(slide, num, 1.24, y + 0.17, 0.5, 0.26, 10.5, accent, True, font=FONT_NUM)
        add_text(slide, title, 1.92, y + 0.12, 3.6, 0.33, 16.5, INK, True)
        add_text(slide, body, 6.2, y + 0.13, 5.4, 0.32, 15, MUTED)
    add_note(slide, extract(source, "### 980元卡的五大成交理由", "### 高定家居精华的五大"))

    # 21 Homecare values
    page += 1
    slide = new_slide(prs, page, "高定家居", "一份赠品，同时解决五个经营问题")
    values = [
        ("转介绍", "让客户有做人情的理由"),
        ("成交", "成为员工980提单工具"),
        ("家居", "打开水乳霜延伸市场"),
        ("到店", "12套按月配置锁频次"),
        ("需求", "每月观察并开发新项目"),
    ]
    for i, (title, body) in enumerate(values):
        x = 0.9 + i * 2.48
        color = [ORANGE, YELLOW, GREEN, BLUE, INK][i]
        add_circle(slide, x + 0.52, 2.55, 1.2, color)
        add_text(slide, f"0{i+1}", x + 0.77, 2.91, 0.7, 0.3, 11, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
        add_text(slide, title, x, 4.05, 2.25, 0.42, 19, INK, True, PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.08, 4.68, 2.08, 0.9, 13.5, MUTED, False, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 高定家居精华的五大经营价值", "## 第五部分"))

    # 22 Card logic
    page += 1
    slide = new_slide(prs, page, "卡项逻辑", "第一张卡，就决定了后面能不能持续升单")
    add_rect(slide, 0.95, 2.4, 5.05, 3.52, PALE_ORANGE, PALE_ORANGE)
    add_text(slide, "当下痛点", 1.35, 2.87, 4.25, 0.52, 26, ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, "解决眼前问题\n建立信任\n促成当下成交", 1.35, 3.75, 4.25, 1.42, 19, INK, False, PP_ALIGN.CENTER)
    add_arrow(slide, 6.22, 3.52, 0.9, 0.68, INK)
    add_rect(slide, 7.35, 2.4, 5.05, 3.52, PALE_GREEN, PALE_GREEN)
    add_text(slide, "最终需求", 7.75, 2.87, 4.25, 0.52, 26, GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, "得到漂亮脸蛋\n维持长期状态\n形成持续复购", 7.75, 3.75, 4.25, 1.42, 19, INK, False, PP_ALIGN.CENTER)
    add_text(slide, "卡项底层逻辑：解决当下痛点，同时通向最终需求", 2.0, 6.32, 9.3, 0.38, 19, INK, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 卡项设计的底层逻辑", "### 好面相"))

    # 23 Good face
    page += 1
    slide = new_slide(prs, page, "最终交付", "好面相的五个字：净、透、润、紧、亮")
    words = [
        ("净", "毛孔毛囊干净", ORANGE),
        ("透", "打开吸收通道", YELLOW),
        ("润", "补进水与营养", GREEN),
        ("紧", "紧致延缓衰老", BLUE),
        ("亮", "整张脸光亮", INK),
    ]
    for i, (word, body, color) in enumerate(words):
        x = 0.82 + i * 2.5
        add_circle(slide, x + 0.55, 2.4, 1.32, color)
        add_text(slide, word, x + 0.85, 2.72, 0.72, 0.55, 29, WHITE, True, PP_ALIGN.CENTER)
        if i < 4:
            add_arrow(slide, x + 1.92, 2.85, 0.5, 0.38, LINE)
        add_text(slide, body, x, 4.08, 2.25, 0.65, 15, INK, True, PP_ALIGN.CENTER)
    add_rect(slide, 1.65, 5.4, 10.0, 0.78, PALE_YELLOW, PALE_YELLOW)
    add_text(slide, "客户看见护肤消费路径，员工看见销售路径", 1.95, 5.64, 9.4, 0.35, 18, INK, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 好面相的五个字", "### 第一张卡"))

    # 24 Card 980
    page += 1
    slide = new_slide(prs, page, "第一张卡", "980元｜一阶调护：净、透")
    add_rect(slide, 0.95, 2.28, 3.15, 3.82, ORANGE, ORANGE)
    add_text(slide, "980", 1.28, 2.72, 2.5, 0.92, 48, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "新客信任卡", 1.28, 3.72, 2.5, 0.42, 19, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "留客 + 裂变", 1.28, 4.58, 2.5, 0.42, 17, "FBDDD5", True, PP_ALIGN.CENTER)
    add_text(slide, "终身限购1次", 1.28, 5.24, 2.5, 0.35, 13.5, WHITE, False, PP_ALIGN.CENTER)
    stages = [
        ("前3次", "清洁毛囊"),
        ("中3次", "消炎修复"),
        ("后3次", "毛孔/闭口修复"),
    ]
    for i, (when, action) in enumerate(stages):
        y = 2.38 + i * 1.08
        add_rect(slide, 4.65, y, 3.25, 0.84, WHITE, LINE)
        add_text(slide, when, 4.93, y + 0.2, 0.95, 0.32, 14, ORANGE, True)
        add_text(slide, action, 6.0, y + 0.18, 1.55, 0.35, 16, INK, True)
    add_rect(slide, 8.45, 2.38, 3.78, 3.1, PALE_YELLOW, PALE_YELLOW)
    add_text(slide, "9 + 9 + 1", 8.83, 2.83, 3.0, 0.68, 34, INK, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "毛囊SPA + 水置换\n+ 白雪公主", 8.83, 3.75, 3.0, 0.92, 16.5, MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, "白雪公主 = 3980升单钩子", 5.0, 5.86, 6.95, 0.42, 17, ORANGE, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 第一张卡：", "### 第二张卡"))

    # 25 Card 3980
    page += 1
    slide = new_slide(prs, page, "第二张卡", "3980元｜二阶养护：净、透、润")
    add_rect(slide, 0.95, 2.28, 3.15, 3.82, GREEN, GREEN)
    add_text(slide, "3980", 1.18, 2.72, 2.7, 0.92, 45, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "二阶养护卡", 1.28, 3.72, 2.5, 0.42, 19, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "5个品种", 1.28, 4.58, 2.5, 0.42, 17, "DCEAE6", True, PP_ALIGN.CENTER)
    add_text(slide, "40 + 3 = 43项", 1.28, 5.24, 2.5, 0.35, 14, WHITE, False, PP_ALIGN.CENTER, font=FONT_NUM)
    combos = [
        ("组合A", "毛囊SPA + 水置换 + 白雪公主"),
        ("组合B", "毛囊SPA + 轮廓塑形 + 颈部淋巴"),
        ("组合C", "白雪公主 + 胶原充盈"),
    ]
    for i, (title, body) in enumerate(combos):
        y = 2.35 + i * 1.08
        add_rect(slide, 4.65, y, 7.55, 0.84, WHITE, LINE)
        add_text(slide, title, 4.95, y + 0.2, 1.05, 0.32, 13.5, GREEN, True)
        add_text(slide, body, 6.15, y + 0.18, 5.65, 0.35, 15.5, INK, True)
    add_rect(slide, 4.65, 5.77, 7.55, 0.48, PALE_GREEN, PALE_GREEN)
    add_text(slide, "10–15次消耗完｜卖疗程方案，不卖项目批发", 4.9, 5.89, 7.05, 0.28, 15.5, GREEN, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 第二张卡：", "### 第三张卡"))

    # 26 Card 6980
    page += 1
    slide = new_slide(prs, page, "第三张卡", "6980元｜三阶综合：净、透、润、紧")
    add_rect(slide, 0.95, 2.28, 3.15, 3.82, BLUE, BLUE)
    add_text(slide, "6980", 1.18, 2.72, 2.7, 0.92, 45, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "三阶综合卡", 1.28, 3.72, 2.5, 0.42, 19, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "40 + 6 = 46项", 1.28, 4.58, 2.5, 0.42, 16, "DFEAF2", True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "+ 4款家居产品", 1.28, 5.24, 2.5, 0.35, 14, WHITE, False, PP_ALIGN.CENTER)
    add_card(slide, 4.62, 2.33, 3.53, 1.48, "轮廓打造", "从皮肤状态进入面部轮廓管理", PALE_BLUE, BLUE, 18, 14)
    add_card(slide, 8.55, 2.33, 3.53, 1.48, "高阶胶原", "导入仪操作，作为进阶体验", PALE_ORANGE, ORANGE, 18, 14)
    add_card(slide, 4.62, 4.16, 3.53, 1.48, "家居复购", "4款家居承接月月复购目标", PALE_GREEN, GREEN, 18, 14)
    add_card(slide, 8.55, 4.16, 3.53, 1.48, "需求培养", "为胶原、水光等需求铺垫", PALE_YELLOW, YELLOW, 18, 14)
    add_note(slide, extract(source, "### 第三张卡：", "## 第六部分"))

    # 27 Annual journey
    page += 1
    slide = new_slide(prs, page, "年消费路径", "三张卡走完，一个客户自然成为“有效客户”")
    journey = [
        ("980", "净 · 透", ORANGE, PALE_ORANGE),
        ("3980", "净 · 透 · 润", GREEN, PALE_GREEN),
        ("6980", "净 · 透 · 润 · 紧", BLUE, PALE_BLUE),
    ]
    for i, (price, result, accent, fill) in enumerate(journey):
        x = 0.95 + i * 4.17
        add_rect(slide, x, 2.55, 3.55, 2.42, fill, fill)
        add_text(slide, price, x + 0.35, 2.92, 2.85, 0.75, 36, accent, True, PP_ALIGN.CENTER, font=FONT_NUM)
        add_text(slide, result, x + 0.35, 3.9, 2.85, 0.42, 17, INK, True, PP_ALIGN.CENTER)
        if i < 2:
            add_arrow(slide, x + 3.65, 3.35, 0.45, 0.42, accent)
    add_rich_text(slide, [("总消费  ", 16, MUTED, False), ("11,840元", 35, ORANGE, True), ("  = 年复购3次 + 年消费1万元以上", 16, INK, True)], 1.5, 5.6, 10.35, 0.7, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "一个客户一年依次购买", "---"))

    # 28 Employee incentive
    page += 1
    slide = new_slide(prs, page, "员工动力", "三级现金积分，把客户升级与员工收益绑定")
    stages = [
        ("成交980", "+300积分", "客户绑定员工", ORANGE),
        ("升单3980", "300提现\n+500积分", "继续维护", GREEN),
        ("升单6980", "500提现\n+500现金", "完成全链路", BLUE),
    ]
    for i, (title, reward, body, accent) in enumerate(stages):
        x = 0.95 + i * 4.17
        add_rect(slide, x, 2.42, 3.55, 3.15, WHITE, LINE)
        add_rect(slide, x, 2.42, 3.55, 0.17, accent)
        add_text(slide, title, x + 0.32, 2.88, 2.9, 0.44, 19, INK, True, PP_ALIGN.CENTER)
        add_text(slide, reward, x + 0.32, 3.62, 2.9, 0.95, 26, accent, True, PP_ALIGN.CENTER, font=FONT_NUM)
        add_text(slide, body, x + 0.32, 4.82, 2.9, 0.34, 13.5, MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, "1个客户全链路 = 1300元｜30个客户 = 39,000元", 1.65, 6.05, 10.0, 0.42, 20, INK, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_note(slide, extract(source, "### 员工三级现金积分奖励", "### 第一次下店"))

    # 29 Three visits
    page += 1
    slide = new_slide(prs, page, "三次下店", "老师按三张卡分阶段培训与驻店带练")
    visits = [
        ("第一次", "30个980种子客户", "交付能力\n老客裂变话术\n新客成交五部曲", ORANGE, PALE_ORANGE),
        ("第二次", "20个3980升单", "980→3980能力\n家居搭配\n补强弱项", GREEN, PALE_GREEN),
        ("第三次", "10个6980升单", "3980→6980能力\n家居销售\n强化全链路", BLUE, PALE_BLUE),
    ]
    for i, (title, goal, body, accent, fill) in enumerate(visits):
        x = 0.95 + i * 4.17
        add_rect(slide, x, 2.35, 3.55, 3.7, fill, fill)
        add_text(slide, title, x + 0.32, 2.72, 2.9, 0.4, 16, accent, True, PP_ALIGN.CENTER)
        add_text(slide, goal, x + 0.32, 3.3, 2.9, 0.52, 21, INK, True, PP_ALIGN.CENTER)
        add_line(slide, x + 0.42, 4.0, x + 3.13, 4.0, "CFC8BD", 0.8)
        add_text(slide, body, x + 0.42, 4.3, 2.7, 1.25, 15, MUTED, False, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 第一次下店", "## 第七部分"))

    # 30 Cooperation
    page += 1
    slide = new_slide(prs, page, "合作方式", "3980元对赌 + 4万元铺货 + 按单补货")
    flow = [
        ("01", "支付3980", "签六个月对赌", ORANGE),
        ("02", "铺货4万", "产品/展示/物料", YELLOW),
        ("03", "小程序下单", "自动分账", GREEN),
        ("04", "卖一补一", "无需持续囤货", BLUE),
    ]
    for i, (num, title, body, accent) in enumerate(flow):
        x = 0.75 + i * 3.12
        add_circle(slide, x + 0.9, 2.45, 1.2, accent)
        add_text(slide, num, x + 1.16, 2.81, 0.7, 0.3, 11, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
        if i < 3:
            add_arrow(slide, x + 2.2, 2.86, 0.73, 0.42, LINE)
        add_text(slide, title, x + 0.25, 4.02, 2.5, 0.42, 18, INK, True, PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.25, 4.6, 2.5, 0.65, 14, MUTED, False, PP_ALIGN.CENTER)
    add_rect(slide, 1.65, 5.72, 10.0, 0.64, PALE_ORANGE, PALE_ORANGE)
    add_text(slide, "硬性对赌：六个月新增100个成交980元的新客", 1.95, 5.91, 9.4, 0.32, 17.5, ORANGE, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 3980元对赌金", "### 分账比例"))

    # 31 Split
    page += 1
    slide = new_slide(prs, page, "分账规则", "门店拿什么，公司承担什么")
    add_rect(slide, 0.95, 2.28, 3.55, 3.9, PALE_ORANGE, PALE_ORANGE)
    add_text(slide, "980", 1.3, 2.68, 2.85, 0.7, 36, ORANGE, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "门店50%\n公司50%", 1.3, 3.62, 2.85, 0.95, 21, INK, True, PP_ALIGN.CENTER)
    add_text(slide, "公司承担产品\n2套 + 12套高定精华", 1.3, 4.98, 2.85, 0.76, 14, MUTED, False, PP_ALIGN.CENTER)
    add_rect(slide, 4.88, 2.28, 3.55, 3.9, PALE_GREEN, PALE_GREEN)
    add_text(slide, "3980 / 6980", 5.18, 2.68, 2.95, 0.7, 31, GREEN, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "门店65%\n公司35%", 5.18, 3.62, 2.95, 0.95, 21, INK, True, PP_ALIGN.CENTER)
    add_text(slide, "公司承担产品、赠品\n员工与老师奖励", 5.18, 4.98, 2.95, 0.76, 14, MUTED, False, PP_ALIGN.CENTER)
    add_rect(slide, 8.81, 2.28, 3.55, 3.9, PALE_BLUE, PALE_BLUE)
    add_text(slide, "前10张980", 9.11, 2.68, 2.95, 0.7, 31, BLUE, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "不分账\n公司收4900", 9.11, 3.62, 2.95, 0.95, 21, INK, True, PP_ALIGN.CENTER)
    add_text(slide, "作为4万元周转货物\n前期回收费用", 9.11, 4.98, 2.95, 0.76, 14, MUTED, False, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 分账比例", "## 第八部分"))

    # 32 Q&A first half
    page += 1
    slide = new_slide(prs, page, "现场Q&A", "先处理钱、客户定义和对赌边界")
    qas = [
        ("原来业绩也分？", "只分小程序成交的三张卡"),
        ("3980为什么还分？", "3980是对赌金，不是加盟费"),
        ("什么算新增新客？", "老客不算；老客转来的B、C算"),
        ("3980/6980怎么分？", "门店65%，公司35%"),
        ("三张卡各做100个？", "只硬性对赌100个新增980"),
    ]
    for i, (q, a) in enumerate(qas):
        y = 2.2 + i * 0.86
        add_text(slide, f"Q{i+1}", 0.98, y + 0.15, 0.52, 0.3, 10.5, ORANGE, True, font=FONT_NUM)
        add_text(slide, q, 1.72, y + 0.1, 3.55, 0.36, 15.5, INK, True)
        add_text(slide, a, 5.72, y + 0.1, 6.0, 0.36, 15, MUTED)
        add_line(slide, 0.98, y + 0.66, 12.18, y + 0.66, LINE, 0.7)
    add_note(slide, extract(source, "### Q1：", "### Q6："))

    # 33 Q&A second half
    page += 1
    slide = new_slide(prs, page, "现场Q&A", "再处理执行、利润、产能与业绩预期")
    qas = [
        ("老师只培训？", "会带练、做客户、参与成交，但核心是教会员工"),
        ("第一批客户哪来？", "先做老客：一是练兵，二是裂变"),
        ("980利润多少？", "主播算法：490 - 180手工费 = 310元"),
        ("980太多做不过来？", "终身限购1次，重点在3980和6980"),
        ("老师能做多少业绩？", "先盘床位与五项数据，不是来做一场活动"),
        ("公司赚什么？", "主播称只赚5个点；1位教练最多20家店"),
    ]
    for i, (q, a) in enumerate(qas):
        y = 2.12 + i * 0.72
        add_text(slide, f"Q{i+6}", 0.98, y + 0.12, 0.55, 0.28, 10.5, GREEN, True, font=FONT_NUM)
        add_text(slide, q, 1.72, y + 0.07, 3.35, 0.34, 14.5, INK, True)
        add_text(slide, a, 5.45, y + 0.07, 6.35, 0.34, 14, MUTED)
        add_line(slide, 0.98, y + 0.55, 12.18, y + 0.55, LINE, 0.7)
    add_note(slide, extract(source, "### Q6：", "## 第九部分"))

    # 34 Close
    page += 1
    slide = new_slide(prs, page, "成交收口", "再次总结系统，再把动作压缩到“领取合同”")
    recap = [
        ("三张卡", "980 / 3980 / 6980"),
        ("两条路径", "公域引流 / 私域裂变"),
        ("五项数据", "客户 / 到店 / 项目 / 价格 / 复购"),
        ("六个月", "培训 + 驻店 + 线上陪跑"),
    ]
    for i, (title, body) in enumerate(recap):
        x = 0.82 + i * 3.08
        add_rect(slide, x, 2.35, 2.75, 1.45, [PALE_ORANGE, PALE_GREEN, PALE_BLUE, PALE_YELLOW][i], LINE)
        add_text(slide, title, x + 0.2, 2.68, 2.35, 0.4, 18, INK, True, PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.2, 3.15, 2.35, 0.3, 12.5, MUTED, False, PP_ALIGN.CENTER)
    add_arrow(slide, 5.65, 4.25, 2.0, 0.72, ORANGE)
    add_rect(slide, 3.75, 5.25, 5.82, 0.92, INK)
    add_text(slide, "向李伟老师领取3980元对赌合同", 4.02, 5.51, 5.28, 0.4, 19, WHITE, True, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "## 第九部分", None))

    # 35 End
    page += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, INK)
    add_rect(slide, 0, 0, 0.22, 7.5, ORANGE)
    add_text(slide, "全场最后一句", 0.92, 1.18, 2.5, 0.38, 12, YELLOW, True)
    add_text(slide, "客户请客，\n我买单。", 0.88, 1.95, 8.2, 1.85, 42, WHITE, True)
    add_line(slide, 0.92, 4.35, 8.85, 4.35, "4C4C4E", 1)
    add_text(slide, "这句话同时完成：利益解释、转介绍降压、裂变机制记忆", 0.92, 4.75, 8.3, 0.48, 18, "D8D4CD")
    add_rect(slide, 9.58, 0, 3.75, 7.5, ORANGE)
    add_text(slide, "END", 10.15, 2.75, 2.6, 0.55, 22, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "35页", 10.15, 3.55, 2.6, 0.55, 19, WHITE, True, PP_ALIGN.CENTER, font=FONT_NUM)
    add_text(slide, "原话详见\n演讲者备注", 10.15, 4.4, 2.6, 0.78, 15, "FBDDD5", False, PP_ALIGN.CENTER)
    add_note(slide, extract(source, "### 结束前再次补讲裂变", None))

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return len(prs.slides)


if __name__ == "__main__":
    count = build_deck()
    print(f"created {OUTPUT} with {count} slides")
